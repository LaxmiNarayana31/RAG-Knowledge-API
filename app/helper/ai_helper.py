from PyPDF2 import PdfReader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain.memory import ConversationBufferMemory
from langchain.chains.conversational_retrieval.base import ConversationalRetrievalChain
from langchain_core.documents import Document
from langchain.prompts import ChatPromptTemplate
from app.helper.llm_helper import LlmHelper, GoogleGeminiEmbeddings
import os
import pandas as pd
import tempfile
import docx
import pptx
from odf.opendocument import load
from odf.text import P
from striprtf.striprtf import rtf_to_text
from spire.presentation import Presentation, FileFormat
from spire.doc import Document as SpireDocument
from spire.doc import FileFormat as SpireFileFormat
from spire.xls import Workbook as SpireWorkbook
from spire.xls import ExcelVersion
from pathlib import Path
from app.utils.exception_handler import handle_exception


class AIHelper:
    vectorstore_folder = "vector_index"

    # Extract text from different document formats
    def extract_text(file):
        '''
        Extract text from different document formats
        Supported formats: .pdf, .doc, .docx, .txt, .ppt, .pptx, .odt, .rtf, .csv, .xls, .xlsx
        '''
        try:
            file_name = file.name.lower()
            text = ""

            match file_name.split('.')[-1]:
                case "pdf":
                    pdf_reader = PdfReader(file)
                    for page in pdf_reader.pages:
                        page_text = page.extract_text()
                        if page_text:
                            text += page_text
                case "doc":
                    # Create a temporary file to save the DOC content
                    with tempfile.NamedTemporaryFile(delete = False, suffix = ".doc") as tmp_file:
                        tmp_file.write(file.read())
                        tmp_file_path = tmp_file.name

                    # Convert DOC to DOCX using Spire
                    document = SpireDocument()
                    document.LoadFromFile(tmp_file_path)
                    document.SaveToFile(tmp_file_path, SpireFileFormat.Docx2019)
                    document.Close()

                    doc = docx.Document(tmp_file_path)
                    text = "\n".join([para.text for para in doc.paragraphs])
                case "docx":
                    doc = docx.Document(file)
                    text = "\n".join([para.text for para in doc.paragraphs])
                case "txt":
                    text = file.read().decode("utf-8")
                case "ppt":
                    # Create a temporary file to save the PPT content
                    with tempfile.NamedTemporaryFile(delete = False, suffix = ".ppt") as tmp_file:
                        tmp_file.write(file.read())
                        tmp_file_path = tmp_file.name

                    # Convert PPT to PPTX using Spire
                    pre = Presentation()
                    pre.LoadFromFile(tmp_file_path)
                    pre.SaveToFile(tmp_file_path, FileFormat.Pptx2019)
                    pre.Dispose()

                    # Load the converted PPTX file
                    presentation = pptx.Presentation(tmp_file_path)
                    for slide in presentation.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text += shape.text + "\n"
                    os.remove(tmp_file_path)
                case "pptx":
                    presentation = pptx.Presentation(file)
                    for slide in presentation.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text += shape.text + "\n"
                case "odt":
                    odt_doc = load(file)
                    # Extract all paragraphs from the ODT file
                    all_paragraphs = odt_doc.getElementsByType(P)
                    for paragraph in all_paragraphs:
                        # Extract text content from each paragraph
                        paragraph_text = ""
                        for node in paragraph.childNodes:
                            if node.nodeType == node.TEXT_NODE:
                                paragraph_text += node.data
                            elif node.nodeType == node.ELEMENT_NODE and node.tagName == "text:span":
                                # Handle nested text spans
                                for child_node in node.childNodes:
                                    if child_node.nodeType == child_node.TEXT_NODE:
                                        paragraph_text += child_node.data
                        text += paragraph_text + "\n"
                case "rtf":
                    rtf_content = file.read()
                    text = rtf_to_text(rtf_content.decode("latin-1"))
                case "csv":
                    df = pd.read_csv(file)
                    rows = []
                    for _, row in df.iterrows():
                        # row_text = ", ".join(f"{col}: {row[col]}" for col in df.columns)
                        row_text = " | ".join(f"{col}: {row[col]}" for col in df.columns)
                        rows.append(row_text)

                    text = "\n".join(rows)
                case "xls":
                    with tempfile.NamedTemporaryFile(delete=False, suffix=".xls") as tmp_file:
                        tmp_file.write(file.read())
                        tmp_file_path = tmp_file.name

                    # Convert XLS to XLSX using Spire
                    workbook = SpireWorkbook()
                    workbook.LoadFromFile(tmp_file_path)
                    workbook.SaveToFile(tmp_file_path, ExcelVersion.Version2016)
                    workbook.Dispose()

                    df = pd.read_excel(tmp_file_path)
                    text = "\n".join([", ".join([f"{col}: {str(row[col])}" for col in df.columns]) for _, row in df.iterrows()])
                case "xlsx":
                    df = pd.read_excel(file)
                    text = "\n".join([", ".join([f"{col}: {str(row[col])}" for col in df.columns]) for _, row in df.iterrows()])
                case _:
                    print(f"Unsupported file format: {file_name}")

            return text
        except Exception as e:
            return handle_exception(e)


    # Create or load a vectorstore from documents
    def build_vectorstore(document_files, document_uuid=None):
        try:
            documents = []
            for file in document_files:
                file_name = file.name.lower()
                text = AIHelper.extract_text(file)
                
                if text:
                    text_splitter = RecursiveCharacterTextSplitter(
                        separators="\n",
                        chunk_size=1000,
                        chunk_overlap=200,
                        length_function=len
                    )
                    chunks = text_splitter.split_text(text)

                    for chunk in chunks:
                        metadata = {"document_name": file_name}
                        if document_uuid:
                            metadata["document_uuid"] = document_uuid
                        document = Document(
                            page_content=chunk,
                            metadata=metadata
                        )
                        documents.append(document)
                            
            embeddings = GoogleGeminiEmbeddings()

            vectorstore_path = Path(AIHelper.vectorstore_folder)

            if vectorstore_path.exists():
                vectorstore = FAISS.load_local(AIHelper.vectorstore_folder, embeddings, allow_dangerous_deserialization = True)
                vectorstore.add_documents(documents)
            else:
                vectorstore = FAISS.from_documents(documents, embedding = embeddings)

            vectorstore.save_local(AIHelper.vectorstore_folder)

            return vectorstore
        except Exception as e:
            return handle_exception(e)


    # Create a conversation chain for conversational retrieval
    def initialize_conversation_chain(vectorstore):
        try:
            if vectorstore is None:
                return "No vectorstore available. Please upload a valid PDF file with text content."

            memory = ConversationBufferMemory(memory_key = 'chat_history', return_messages = True)
            retriever = vectorstore.as_retriever(search_type = "mmr",  search_kwargs = {"k": 5, "fetch_k": 10})

            if retriever is None:
                return "Error: Could not create retriever. Ensure vectorstore is valid and properly initialized."
            
            llm = LlmHelper.googleGeminiLlm()
            conversation_chain = ConversationalRetrievalChain.from_llm(
                llm = llm,
                retriever = retriever,
                memory = memory
            )
            return conversation_chain
        except Exception as e:
            return handle_exception(e)
    
    
    # Get response for user question from the llm
    def get_llm_response(user_question):
        try:
            print("================= user_question ================", user_question)
            
            embeddings = GoogleGeminiEmbeddings()
            vectorstore_path = Path(AIHelper.vectorstore_folder)
            
            if not vectorstore_path.exists():
                return "Vectorstore is not available. Please upload a document first."

            vectorstore = FAISS.load_local(folder_path = AIHelper.vectorstore_folder, embeddings = embeddings, allow_dangerous_deserialization = True)
            
            conversation_chain = AIHelper.initialize_conversation_chain(vectorstore)
            if isinstance(conversation_chain, str):
                return conversation_chain
            
            prompt_template = ChatPromptTemplate.from_template(
                """You are an advanced AI assistant answering questions strictly based on the provided document.  

                Conversation History: {chat_history}

                User's Question: "{question}"  

                Relevant Document Sections: {context}  

                ### Instructions:  
                1. If the question is unrelated to the document, say: "The question is out of context."  
                2. Answer based strictly on the document, ensuring accuracy.  
                3. Provide a structured response, explaining your reasoning when necessary.  
                4. If unsure, state: "I do not have enough information from the document."  
                """
            )

            formatted_question = prompt_template.format(question = user_question, context = "", chat_history = "")
            chat_history = None
            input_dict = {'question': formatted_question}
            if chat_history:
                input_dict['chat_history'] = chat_history

            response = conversation_chain(input_dict)
            llm_response = response['chat_history'][-1].content.replace("\n", " ")
            print("================= llm_response ================", llm_response)
            return llm_response
        except Exception as e:
            return handle_exception(e)

